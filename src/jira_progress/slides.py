"""PPTX exporter for progress reports — Red Hat brand-aligned.

Slide structure:
  - Outcome entry: 1 overview slide + 1 detail slide per STRAT
  - STRAT entry: 1 detail slide per STRAT

Visual language: red (#EE0000) + black (#292929) only, except health dots.
Grey (#F5F5F5) section blocks. Red accent line. Proper footer.
"""

from __future__ import annotations

import json
import logging
import re
from datetime import date as _date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

BRAND_PATH = Path(__file__).resolve().parent.parent / "config" / "brand.json"
PROGRESS_OUTPUT_DIR = Path("progress")

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN_L = Inches(0.7)
MARGIN_R = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
MAX_Y = Inches(5.0)

RED = RGBColor(0xEE, 0x00, 0x00)
BLACK = RGBColor(0x29, 0x29, 0x29)
GREY_LIGHT = RGBColor(0x4D, 0x4D, 0x4D)
GREY_BG = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_HEALTH_DOT = {"green": "●", "yellow": "●", "red": "●", "unknown": "○"}
HEADING_FONT = "Red Hat Display"
BODY_FONT = "Red Hat Text"
JIRA_BROWSE = "https://issues.redhat.com/browse/"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _load_brand() -> dict:
    return json.loads(BRAND_PATH.read_text())


def _hex(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _status_to_category(status: str) -> str:
    s = status.lower()
    if s in ("closed", "done", "resolved", "release pending", "released"):
        return "Done"
    if s in ("in progress", "in review", "review", "code review", "testing"):
        return "In Progress"
    return "To Do"


def _trunc(text: str, limit: int) -> str:
    if not text:
        return ""
    if len(text) > limit:
        return text[: limit - 3] + "..."
    return text


def _release_sort_key(release: str) -> tuple:
    """Sort: EA milestones before GA, Unplanned/Other last."""
    if release in ("Unplanned", "Other"):
        return (999, 999, 999, 0)
    m = re.match(r"(\d+)\.(\d+)(?:-(ea)-?(\d+))?", release, re.IGNORECASE)
    if not m:
        return (998, 0, 0, 0)
    major, minor = int(m.group(1)), int(m.group(2))
    if m.group(3):
        return (major, minor, 0, int(m.group(4)))
    return (major, minor, 1, 0)


def _sorted_releases(strat_by_release: dict) -> list[str]:
    return sorted(strat_by_release.keys(), key=_release_sort_key)


def _release_label(release: str) -> str:
    if release in ("Unplanned", "Other"):
        return release
    m = re.match(r"(\d+\.\d+)(?:-(ea)-?(\d+))?", release, re.IGNORECASE)
    if not m:
        return release
    if m.group(2):
        return f"{release} (Early Access)"
    return f"{release} (GA)"


def _health_counts(strats) -> str:
    counts: dict[str, int] = {}
    for s in strats:
        counts[s.health] = counts.get(s.health, 0) + 1
    parts = []
    for h in ("green", "yellow", "red", "unknown"):
        if counts.get(h):
            parts.append(f"{counts[h]} {h.upper()}")
    return ", ".join(parts) if parts else "—"


def _impl_counts_for_strat(strat_impl: list[dict]) -> dict[str, dict[str, int]]:
    by_project: dict[str, dict[str, int]] = {}
    for t in strat_impl:
        proj = t.get("project", "?")
        cat = _status_to_category(t.get("status", ""))
        by_project.setdefault(proj, {"Done": 0, "In Progress": 0, "To Do": 0})
        by_project[proj][cat] = by_project[proj].get(cat, 0) + 1
    return by_project


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------


def _add_text(
    slide, left, top, width, height, text: str, *,
    size: int = 10, bold: bool = False, italic: bool = False,
    font: str = BODY_FONT, color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT, wrap: bool = True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    return box


def _add_jira_link_run(paragraph, key: str, *, size: int = 8, bold: bool = False,
                       font: str = BODY_FONT, color: RGBColor = BLACK):
    """Add a hyperlinked run for a Jira issue key."""
    r = paragraph.add_run()
    r.text = key
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    r.hyperlink.address = f"{JIRA_BROWSE}{key}"
    return r


def _add_grey_block(slide, left, top, width, height):
    """Sharp-cornered grey background for a section (no rounding, no shadow)."""
    from pptx.oxml.ns import qn
    block = slide.shapes.add_shape(5, left, top, width, height)
    block.fill.solid()
    block.fill.fore_color.rgb = GREY_BG
    block.line.fill.background()
    sp_pr = block._element.spPr
    prstGeom = sp_pr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "rect")
    effect_lst = sp_pr.find(qn("a:effectLst"))
    if effect_lst is not None:
        sp_pr.remove(effect_lst)
    return block


# ---------------------------------------------------------------------------
# Slide chrome (brand template elements)
# ---------------------------------------------------------------------------


def _add_slide_chrome(slide, section_text: str, page_num: int) -> None:
    """Add Red Hat brand chrome: accent line, section marker, footer."""
    # Red vertical accent line (left edge)
    accent = slide.shapes.add_connector(
        1, Inches(0.25), Inches(0.2), Inches(0.25), Inches(0.65),
    )
    accent.line.width = Pt(3)
    accent.line.color.rgb = RED

    # Section marker (red text, top-left) — hyperlinked to Jira
    marker_box = slide.shapes.add_textbox(MARGIN_L, Inches(0.2), Inches(4), Inches(0.25))
    marker_tf = marker_box.text_frame
    marker_tf.word_wrap = False
    _add_jira_link_run(marker_tf.paragraphs[0], section_text, size=9, color=RED)

    # Footer: page number (bottom-left)
    _add_text(slide, Inches(0.25), Inches(5.15), Inches(0.4), Inches(0.2),
              str(page_num), size=7, font=BODY_FONT, color=GREY_LIGHT)

    # Footer: source (bottom-left)
    _add_text(slide, MARGIN_L, Inches(5.15), Inches(3), Inches(0.35),
              f"Source: Jira\n{_date.today().isoformat()}",
              size=6, font=BODY_FONT, color=GREY_LIGHT)

    # Footer: Red Hat hat placeholder (bottom-right) — a small red circle
    hat = slide.shapes.add_shape(
        9, Inches(9.3), Inches(5.15), Inches(0.25), Inches(0.25),
    )
    hat.fill.solid()
    hat.fill.fore_color.rgb = RED
    hat.line.fill.background()


# ---------------------------------------------------------------------------
# Overview slide
# ---------------------------------------------------------------------------


def add_overview_slide(prs: Presentation, data, brand: dict, page_num: int = 1) -> None:
    """Overview: release-grouped STRAT list, impl counts, deps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hc = brand.get("health_colors", {})

    _add_slide_chrome(slide, data.outcome_key, page_num)

    # Title (2 lines max)
    _add_text(slide, MARGIN_L, Inches(0.5), Inches(8.5), Inches(0.7),
              _trunc(data.outcome_summary, 120),
              size=22, bold=True, font=HEADING_FONT, color=BLACK)

    # Status line
    total_strats = len(data.strat_health)
    _add_text(slide, MARGIN_L, Inches(1.15), Inches(8.5), Inches(0.25),
              f"Status: {data.outcome_status}  ·  {total_strats} STRATs  ·  {_health_counts(data.strat_health)}",
              size=10, font=BODY_FONT, color=GREY_LIGHT)

    # --- Release-grouped STRAT list (left ~60%) ---
    release_order = _sorted_releases(data.strat_by_release)
    list_x = MARGIN_L
    list_w = Inches(5.8)
    list_y = Inches(1.55)

    list_box = slide.shapes.add_textbox(list_x, list_y, list_w, MAX_Y - list_y)
    ltf = list_box.text_frame
    ltf.word_wrap = True
    first_para = True

    for release in release_order:
        strats = data.strat_by_release.get(release, [])
        p = ltf.paragraphs[0] if first_para else ltf.add_paragraph()
        first_para = False
        p.space_before = Pt(6)
        p.space_after = Pt(2)

        rh = p.add_run()
        rh.text = _release_label(release)
        rh.font.size = Pt(10)
        rh.font.bold = True
        rh.font.name = HEADING_FONT
        rh.font.color.rgb = RED

        rc = p.add_run()
        rc.text = f"  {len(strats)} STRAT{'s' if len(strats) != 1 else ''}"
        rc.font.size = Pt(9)
        rc.font.name = BODY_FONT
        rc.font.color.rgb = GREY_LIGHT

        for sh in strats[:8]:
            dot = _HEALTH_DOT.get(sh.health, "○")
            dot_color = _hex(hc.get(sh.health, hc.get("unknown", "#A3A3A3")))

            ps = ltf.add_paragraph()
            ps.space_after = Pt(1)

            rd = ps.add_run()
            rd.text = f"  {dot} "
            rd.font.size = Pt(9)
            rd.font.bold = True
            rd.font.color.rgb = dot_color
            rd.font.name = BODY_FONT

            _add_jira_link_run(ps, sh.key, size=8, bold=False, color=BLACK)

            rk = ps.add_run()
            rk.text = f" — {_trunc(sh.summary, 65)}"
            rk.font.size = Pt(8)
            rk.font.name = BODY_FONT
            rk.font.color.rgb = BLACK

        if len(strats) > 8:
            po = ltf.add_paragraph()
            po.text = f"    +{len(strats) - 8} more"
            po.font.size = Pt(7)
            po.font.italic = True
            po.font.name = BODY_FONT
            po.font.color.rgb = GREY_LIGHT

    # --- Right panel: impl counts in grey block ---
    panel_x = Inches(6.8)
    panel_w = Inches(2.7)
    panel_y = Inches(1.55)
    panel_h = Inches(2.0)

    _add_grey_block(slide, panel_x, panel_y, panel_w, panel_h)

    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.1), panel_w - Inches(0.3), Inches(0.2),
              "Implementation", size=10, bold=True, font=HEADING_FONT, color=RED)

    total_all = done_all = ip_all = todo_all = 0
    impl_y = panel_y + Inches(0.4)

    for project, cats in sorted(data.impl_counts.items()):
        d = cats.get("Done", 0)
        ip = cats.get("In Progress", 0)
        t = cats.get("To Do", 0)
        total = d + ip + t
        total_all += total
        done_all += d
        ip_all += ip
        todo_all += t

        _add_text(slide, panel_x + Inches(0.15), impl_y, panel_w - Inches(0.3), Inches(0.18),
                  f"{project}  {total} ({d}D / {ip}IP / {t}T)",
                  size=8, font=BODY_FONT, color=BLACK)
        impl_y += Inches(0.2)

    impl_y += Inches(0.05)
    _add_text(slide, panel_x + Inches(0.15), impl_y, panel_w - Inches(0.3), Inches(0.18),
              f"Total: {total_all}  —  {done_all}D / {ip_all}IP / {todo_all}T",
              size=8, bold=True, font=BODY_FONT, color=BLACK)

    # --- Bottom: dependencies ---
    if data.strat_links:
        dep_parts = []
        for sl in data.strat_links[:4]:
            dep_parts.append(f"{sl.source_key} {sl.link_description} {sl.target_key}")
        extra = len(data.strat_links) - 4
        dep_text = " · ".join(dep_parts)
        if extra > 0:
            dep_text += f"  (+{extra} more)"
        _add_text(slide, MARGIN_L, MAX_Y, CONTENT_W, Inches(0.2),
                  _trunc(dep_text, 130), size=7, italic=True, font=BODY_FONT, color=GREY_LIGHT)


# ---------------------------------------------------------------------------
# STRAT detail slide
# ---------------------------------------------------------------------------


def add_strat_slide(prs: Presentation, sh, strat_impl: list[dict], strat_links: list, brand: dict, page_num: int = 1) -> None:
    """Detail slide for a STRAT: grey-blocked sections for signals, impl, relationships."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hc = brand.get("health_colors", {})

    _add_slide_chrome(slide, sh.key, page_num)

    # Title (2 lines max)
    _add_text(slide, MARGIN_L, Inches(0.5), Inches(8.5), Inches(0.7),
              _trunc(sh.summary, 110),
              size=20, bold=True, font=HEADING_FONT, color=BLACK)

    # Status line with health dot
    dot = _HEALTH_DOT.get(sh.health, "○")
    dot_color = _hex(hc.get(sh.health, hc.get("unknown", "#A3A3A3")))
    release_text = sh.target_release or "Unplanned"

    status_box = slide.shapes.add_textbox(MARGIN_L, Inches(1.15), Inches(8.5), Inches(0.25))
    stf = status_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]

    rd = sp.add_run()
    rd.text = f"{dot} "
    rd.font.size = Pt(11)
    rd.font.bold = True
    rd.font.color.rgb = dot_color
    rd.font.name = BODY_FONT

    rs = sp.add_run()
    rs.text = f"{sh.health.upper()}  ·  {sh.status}  ·  {_release_label(release_text)}"
    rs.font.size = Pt(10)
    rs.font.name = BODY_FONT
    rs.font.color.rgb = GREY_LIGHT

    # ---- SECTION: Health Signals (grey block, left) ----
    sec_x = MARGIN_L
    sec_y = Inches(1.55)
    sec_w = Inches(5.3)
    sec_h = Inches(1.8)

    _add_grey_block(slide, sec_x, sec_y, sec_w, sec_h)
    _add_text(slide, sec_x + Inches(0.15), sec_y + Inches(0.08), sec_w - Inches(0.3), Inches(0.2),
              "Health Signals", size=10, bold=True, font=HEADING_FONT, color=RED)

    signal_lines = [
        ("Status", sh.status),
        ("Color Status", sh.color_status or "Not set"),
        ("Status Summary", _trunc(sh.status_summary, 90) or "—"),
        ("Blocked", f"Yes — {_trunc(sh.blocked_reason, 50)}" if sh.blocked else "No"),
        ("Target Release", f"{sh.target_release or '—'} ({sh.release_commitment})"),
    ]

    sig_box = slide.shapes.add_textbox(sec_x + Inches(0.15), sec_y + Inches(0.35), sec_w - Inches(0.3), sec_h - Inches(0.45))
    sig_tf = sig_box.text_frame
    sig_tf.word_wrap = True

    for i, (label, value) in enumerate(signal_lines):
        p = sig_tf.paragraphs[0] if i == 0 else sig_tf.add_paragraph()
        p.space_after = Pt(2)

        rl = p.add_run()
        rl.text = f"{label}:  "
        rl.font.size = Pt(9)
        rl.font.bold = True
        rl.font.name = BODY_FONT
        rl.font.color.rgb = BLACK

        rv = p.add_run()
        rv.text = value
        rv.font.size = Pt(9)
        rv.font.name = BODY_FONT
        rv.font.color.rgb = GREY_LIGHT

    # ---- SECTION: Relationships (grey block, right) ----
    rel_x = Inches(6.3)
    rel_w = Inches(3.2)
    rel_h = Inches(1.8)

    _add_grey_block(slide, rel_x, sec_y, rel_w, rel_h)
    _add_text(slide, rel_x + Inches(0.15), sec_y + Inches(0.08), rel_w - Inches(0.3), Inches(0.2),
              "Relationships", size=10, bold=True, font=HEADING_FONT, color=RED)

    relevant_links = [
        sl for sl in strat_links
        if sl.source_key == sh.key or sl.target_key == sh.key
    ]

    if relevant_links:
        rel_box = slide.shapes.add_textbox(rel_x + Inches(0.15), sec_y + Inches(0.35), rel_w - Inches(0.3), rel_h - Inches(0.45))
        rtf = rel_box.text_frame
        rtf.word_wrap = True

        for i, sl in enumerate(relevant_links[:6]):
            p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
            p.space_after = Pt(2)

            if sl.source_key == sh.key:
                r = p.add_run()
                r.text = f"→ {sl.link_description} "
                r.font.size = Pt(8)
                r.font.name = BODY_FONT
                r.font.color.rgb = BLACK
                _add_jira_link_run(p, sl.target_key, size=8, color=BLACK)
            else:
                r = p.add_run()
                r.text = f"← {sl.link_description} (from "
                r.font.size = Pt(8)
                r.font.name = BODY_FONT
                r.font.color.rgb = BLACK
                _add_jira_link_run(p, sl.source_key, size=8, color=BLACK)
                r2 = p.add_run()
                r2.text = ")"
                r2.font.size = Pt(8)
                r2.font.name = BODY_FONT
                r2.font.color.rgb = BLACK

        if len(relevant_links) > 6:
            po = rtf.add_paragraph()
            po.text = f"+{len(relevant_links) - 6} more"
            po.font.size = Pt(7)
            po.font.italic = True
            po.font.name = BODY_FONT
            po.font.color.rgb = GREY_LIGHT
    else:
        _add_text(slide, rel_x + Inches(0.15), sec_y + Inches(0.4), rel_w - Inches(0.3), Inches(0.2),
                  "No inter-STRAT relationships", size=8, italic=True, font=BODY_FONT, color=GREY_LIGHT)

    # ---- SECTION: Implementation Tickets (grey block, full width bottom) ----
    impl_y = Inches(3.55)
    impl_w = CONTENT_W
    impl_h = Inches(1.35)

    _add_grey_block(slide, sec_x, impl_y, impl_w, impl_h)
    _add_text(slide, sec_x + Inches(0.15), impl_y + Inches(0.08), impl_w - Inches(0.3), Inches(0.2),
              "Implementation Tickets", size=10, bold=True, font=HEADING_FONT, color=RED)

    by_project = _impl_counts_for_strat(strat_impl)
    total_all = len(strat_impl)

    if by_project:
        impl_box = slide.shapes.add_textbox(sec_x + Inches(0.15), impl_y + Inches(0.35), impl_w - Inches(0.3), impl_h - Inches(0.45))
        itf = impl_box.text_frame
        itf.word_wrap = True

        for i, (proj, cats) in enumerate(sorted(by_project.items())):
            d, ip, t = cats.get("Done", 0), cats.get("In Progress", 0), cats.get("To Do", 0)
            p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
            p.space_after = Pt(2)

            rp = p.add_run()
            rp.text = f"{proj}  "
            rp.font.size = Pt(9)
            rp.font.bold = True
            rp.font.name = BODY_FONT
            rp.font.color.rgb = BLACK

            rv = p.add_run()
            rv.text = f"{d + ip + t} total  ({d} Done / {ip} In Progress / {t} To Do)"
            rv.font.size = Pt(9)
            rv.font.name = BODY_FONT
            rv.font.color.rgb = GREY_LIGHT

        pt = itf.add_paragraph()
        pt.space_before = Pt(4)
        rt = pt.add_run()
        done_all = sum(cats.get("Done", 0) for cats in by_project.values())
        ip_all = sum(cats.get("In Progress", 0) for cats in by_project.values())
        rt.text = f"Total: {total_all}  —  {done_all} Done, {ip_all} In Progress"
        rt.font.size = Pt(9)
        rt.font.bold = True
        rt.font.name = BODY_FONT
        rt.font.color.rgb = BLACK
    else:
        _add_text(slide, sec_x + Inches(0.15), impl_y + Inches(0.4), impl_w - Inches(0.3), Inches(0.2),
                  "No implementation tickets found", size=9, italic=True, font=BODY_FONT, color=GREY_LIGHT)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def export_progress_to_pptx(
    data,
    output_dir: Path = PROGRESS_OUTPUT_DIR,
) -> Path:
    """Create a .pptx presentation from ProgressData."""
    brand = _load_brand()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    is_strat_entry = any(
        sh.relationship == "self" and sh.key == data.outcome_key
        for sh in data.strat_health
    )

    page = 1

    if not is_strat_entry:
        add_overview_slide(prs, data, brand, page_num=page)
        page += 1

    for sh in data.strat_health:
        strat_impl = data.strat_impl_map.get(sh.key, [])
        add_strat_slide(prs, sh, strat_impl, data.strat_links, brand, page_num=page)
        page += 1

    output_dir.mkdir(parents=True, exist_ok=True)
    safe_key = data.outcome_key.replace("-", "_").lower()
    filename = f"{_date.today().isoformat()}_{safe_key}.pptx"
    filepath = output_dir / filename
    prs.save(str(filepath))
    logger.info("Saved PPTX to %s", filepath)
    return filepath
